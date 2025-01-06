#!/usr/bin/env python3
"""
Markdown Knowledge Object Factory
-------------------------------
A script to process various document formats into standardized markdown knowledge articles
with intelligent tagging and cross-referencing capabilities.

Can be run directly or used as a tool in an LLM agent toolkit.

Usage as script:
    python markdown_knowledge_object_factory.py [--input-dir DIR] [--output-dir DIR]

Usage as tool:
    from markdown_knowledge_object_factory import MarkdownFactory, process_files
    
    # Direct async usage
    result = await process_files(input_dir="docs", output_dir="knowledge")
    
    # Or use the factory for more control
    factory = MarkdownFactory()
    result = await factory.process_directory("docs", "knowledge")
"""

# Import required libraries
import os
import json
import asyncio
import logging
import argparse
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Any, Optional, Union, Tuple, TypeAlias
from dataclasses import dataclass
from tqdm.asyncio import tqdm

# Type definitions
ProcessingResult: TypeAlias = Dict[str, Any]
ToolResponse: TypeAlias = Tuple[bool, ProcessingResult]

# External dependencies
import chromadb
from chromadb.config import Settings
from dotenv import load_dotenv
from groq import AsyncGroq
from pydantic import BaseModel, Field
import pandas as pd
from pypdf import PdfReader
from docx import Document
import pptx

# Internal dependencies
from .tag_suggester import TagSuggester
from .tagging_system import TaggingRules

# Initialize environment and configuration
load_dotenv()  # Load environment variables from .env file
GROQ_API_KEY = os.getenv("GROQ_API_KEY")  # Get Groq API key for LLM processing

# Set up logging configuration with timestamp and level
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

def log_step(emoji: str, message: str, level: str = "info"):
    """Enhanced logging function that includes emojis for better visual feedback"""
    log_func = getattr(logger, level)
    log_func(f"{emoji} {message}")

class ProcessingConfig(BaseModel):
    """Configuration settings for document processing pipeline"""
    input_dir: str = Field(default="RAG_init")  # Directory containing source files
    output_dir: str = Field(default="RAG_refined")  # Directory for processed markdown files
    supported_extensions: List[str] = Field(  # List of file types we can process
        default=[".pdf", ".md", ".txt", ".docx", ".csv", ".xlsx", ".pptx"]
    )
    chunk_size: int = Field(default=1000)  # Size of text chunks for processing
    chunk_overlap: int = Field(default=200)  # Overlap between chunks for context preservation

class KnowledgeArticle(BaseModel):
    """Data model for processed knowledge articles"""
    title: str  # Article title
    content: str  # Main content in markdown format
    tags: List[str]  # List of contextual tags
    created_date: datetime = Field(default_factory=datetime.now)  # Article creation timestamp
    references: List[str] = Field(default_factory=list)  # Related articles
    source_file: str  # Original source file path
    metadata: Dict[str, Any] = Field(default_factory=dict)  # Additional metadata

class DocumentReader:
    """Handles reading and extracting text from various file formats"""
    
    @staticmethod
    async def read_pdf(file_path: Path) -> str:
        """Extract text from PDF files"""
        reader = PdfReader(str(file_path))
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text

    @staticmethod
    async def read_markdown(file_path: Path) -> str:
        """Read markdown or text files with UTF-8 encoding"""
        return file_path.read_text(encoding='utf-8')

    @staticmethod
    async def read_docx(file_path: Path) -> str:
        """Extract text from Word documents"""
        doc = Document(str(file_path))
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])

    @staticmethod
    async def read_spreadsheet(file_path: Path) -> str:
        """Convert spreadsheet data to markdown table format"""
        df = pd.read_excel(file_path) if file_path.suffix == '.xlsx' else pd.read_csv(file_path)
        return df.to_markdown()

    @staticmethod
    async def read_pptx(file_path: Path) -> str:
        """Extract text from PowerPoint presentations"""
        prs = pptx.Presentation(str(file_path))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

class ContentAnalyzer:
    """Analyzes document content using Groq LLM for metadata generation"""
    
    def __init__(self):
        self.client = AsyncGroq(api_key=GROQ_API_KEY)
        self.model = "mixtral-8x7b-32768"  # Using Mixtral for better analysis
        self.tag_suggester = TagSuggester(self.client, TaggingRules())

    async def analyze_content(self, content: str, source_file: str) -> Dict[str, Any]:
        """
        Enhanced content analysis with sophisticated RPG-style tagging system and knowledge graph metadata.
        Maximum context window: 32,768 tokens for mixtral-8x7b-32768
        """
        # Extract title from content
        title = next((line.strip('# ') for line in content.split('\n') 
                     if line.strip().startswith('# ')), "Untitled Document")
        
        try:
            # Get tag suggestions using TagSuggester
            tag_suggestions = await self.tag_suggester.suggest_tags(content, title, source_file)
            
            # Filter suggestions by confidence
            valid_suggestions = self.tag_suggester.filter_suggestions(tag_suggestions, min_confidence=0.7)
            
            # Organize tags by tier
            tags = {
                "common": {"domain": [], "era": [], "format": []},
                "fine": {"themes": [], "concepts": [], "patterns": []},
                "rare": {"topics": [], "terminology": [], "methods": []},
                "epic": {"insights": [], "connections": [], "innovations": []},
                "legendary": {"principles": [], "paradigms": []}
            }
            
            # Sort suggestions into appropriate tiers and categories
            for suggestion in valid_suggestions:
                tier_name = suggestion.tier.name.lower()
                if tier_name in tags and suggestion.category in tags[tier_name]:
                    tags[tier_name][suggestion.category].append(suggestion.name)
            
            # Format tags with proper emoji indicators and maintain hierarchy
            formatted_tags = []
            
            # Common tags (⚪)
            common_tags = []
            for category, tag_list in tags["common"].items():
                common_tags.extend(tag_list)
            if common_tags:
                formatted_tags.append(f"⚪ {' '.join(['#' + tag.lower().strip() for tag in common_tags])}")
            
            # Fine tags (🟢)
            fine_tags = []
            for category, tag_list in tags["fine"].items():
                fine_tags.extend(tag_list)
            if fine_tags:
                formatted_tags.append(f"🟢 {' '.join(['#' + tag.lower().strip() for tag in fine_tags])}")
            
            # Rare tags (🔵)
            rare_tags = []
            for category, tag_list in tags["rare"].items():
                rare_tags.extend(tag_list)
            if rare_tags:
                formatted_tags.append(f"🔵 {' '.join(['#' + tag.lower().strip() for tag in rare_tags])}")
            
            # Epic tags (🟣)
            epic_tags = []
            for category, tag_list in tags["epic"].items():
                epic_tags.extend(tag_list)
            if epic_tags:
                formatted_tags.append(f"🟣 {' '.join(['#' + tag.lower().strip() for tag in epic_tags])}")
            
            # Legendary tags (🟡)
            legendary_tags = []
            for category, tag_list in tags["legendary"].items():
                legendary_tags.extend(tag_list)
            if legendary_tags:
                formatted_tags.append(f"🟡 {' '.join(['#' + tag.lower().strip() for tag in legendary_tags])}")
            
            # Ensure we have at least some tags
            if not any(formatted_tags):
                formatted_tags = ["⚪ #unprocessed #needs-review"]
            
            # Get related content based on tags
            all_tags = []
            for tier_tags in tags.values():
                for category_tags in tier_tags.values():
                    all_tags.extend(category_tags)
            related_files = self.tag_suggester.get_related_content(all_tags)
            
            # Create the result object
            result = {
                "title": title,
                "file_name": title.lower().replace(' ', '-'),
                "tags": tags,
                "formatted_tags": formatted_tags,
                "metadata": {
                    "tag_confidence": {s.name: s.confidence for s in valid_suggestions},
                    "tag_explanations": {s.name: s.explanation for s in valid_suggestions},
                    "related_files": list(related_files)
                },
                "formatted_content": content
            }
            
            log_step("✨", f"Successfully analyzed content with {len(formatted_tags)} tag groups")
            return result
            
        except Exception as e:
            log_step("⚠️", f"Error analyzing content: {str(e)}", "error")
            return {
                "title": title,
                "file_name": "untitled-document",
                "tags": {"common": {"domain": ["unprocessed"], "era": [], "format": []}},
                "formatted_tags": ["⚪ #unprocessed #needs-review"],
                "metadata": {},
                "formatted_content": content
            }

class KnowledgeBaseManager:
    """Manages the vector database and metadata for knowledge articles"""
    
    def __init__(self, persist_directory: str = ".knowledge_base"):
        # Initialize storage directories
        Path(persist_directory).mkdir(parents=True, exist_ok=True)
        
        # Set up ChromaDB for vector storage
        self.client = chromadb.PersistentClient(path=persist_directory)
        
        # Create or get collection
        try:
            self.collection = self.client.get_or_create_collection(
                name="knowledge_articles",
                metadata={"hnsw:space": "cosine"}
            )
        except Exception as e:
            log_step("⚠️", f"ChromaDB error: {str(e)}", "error")
            # Fallback: Try to create a new collection
            self.collection = self.client.create_collection(
                name="knowledge_articles",
                metadata={"hnsw:space": "cosine"}
            )
            
        self.metadata_file = Path(persist_directory) / "metadata.json"
        
        # Initialize metadata storage
        if not self.metadata_file.exists():
            self.metadata_file.write_text("{}")
            
        self.metadata = self._load_metadata()
        log_step("📚", f"Initialized knowledge base at {persist_directory}")

    def _load_metadata(self) -> Dict:
        """Load metadata from JSON file"""
        if self.metadata_file.exists():
            return json.loads(self.metadata_file.read_text())
        return {}

    def _save_metadata(self):
        """Save metadata to JSON file"""
        self.metadata_file.write_text(json.dumps(self.metadata, indent=2))

    async def add_article(self, article: KnowledgeArticle):
        """Add article to vector store and update metadata"""
        # Prepare metadata for ChromaDB
        metadata = {
            "title": article.title,
            "tags": ", ".join(article.tags),
            "created_date": article.created_date.isoformat()
        }
        
        # Store in vector database
        log_step("💾", f"Adding article to vector store: {article.title}")
        self.collection.add(
            documents=[article.content],
            metadatas=[metadata],
            ids=[article.title]
        )

        # Update metadata store
        self.metadata[article.title] = {
            "tags": article.tags,
            "created_date": article.created_date.isoformat(),
            "source_file": article.source_file
        }
        self._save_metadata()
        log_step("✅", f"Successfully added article: {article.title}")

    async def find_related_articles(self, content: str, n_results: int = 5) -> List[str]:
        """Find semantically similar articles using vector similarity"""
        results = self.collection.query(
            query_texts=[content],
            n_results=n_results
        )
        return [metadata["title"] for metadata in results["metadatas"][0]]

class MarkdownFactory:
    """Main orchestrator for document processing pipeline"""
    
    def __init__(self, config: Optional[ProcessingConfig] = None):
        self.config = config or ProcessingConfig()
        self.document_reader = DocumentReader()
        self.content_analyzer = ContentAnalyzer()
        self.kb_manager = KnowledgeBaseManager()
        self.semaphore = asyncio.Semaphore(5)
        
        # Create necessary directories
        Path(self.config.input_dir).mkdir(exist_ok=True)
        Path(self.config.output_dir).mkdir(exist_ok=True)

    @classmethod
    async def create(cls, input_dir: str = "RAG_init", output_dir: str = "RAG_refined") -> 'MarkdownFactory':
        """Factory method for async initialization"""
        config = ProcessingConfig(input_dir=input_dir, output_dir=output_dir)
        return cls(config)

    async def process_directory(
        self, 
        input_dir: Optional[str] = None, 
        output_dir: Optional[str] = None
    ) -> ToolResponse:
        """Process all supported files in a directory
        
        Args:
            input_dir: Optional override for input directory
            output_dir: Optional override for output directory
            
        Returns:
            Tuple[bool, Dict]: (success, result_data)
            - success: Whether processing completed successfully
            - result_data: Dictionary containing:
                - processed_files: List of successfully processed files
                - failed_files: List of files that failed processing
                - stats: Processing statistics
        """
        try:
            # Update directories if provided
            if input_dir:
                self.config.input_dir = input_dir
            if output_dir:
                self.config.output_dir = output_dir

            # Get list of files to process
            input_path = Path(self.config.input_dir)
            files = [f for f in input_path.glob("*") 
                    if f.suffix.lower() in self.config.supported_extensions]
            
            log_step("🔍", f"Found {len(files)} files to process")
            
            if not files:
                return True, {
                    "processed_files": [],
                    "failed_files": [],
                    "stats": {"total": 0, "successful": 0, "failed": 0}
                }

            # Process files with progress tracking
            tasks = [self.process_file_with_semaphore(file) for file in files]
            results = await tqdm.gather(*tasks, desc="📚 Processing files")

            # Compile results
            successful = [r.source_file for r in results if r is not None]
            failed = [str(f) for f, r in zip(files, results) if r is None]
            
            result_data = {
                "processed_files": successful,
                "failed_files": failed,
                "stats": {
                    "total": len(files),
                    "successful": len(successful),
                    "failed": len(failed)
                }
            }
            
            return True, result_data

        except Exception as e:
            logger.error(f"Error processing directory: {str(e)}")
            return False, {"error": str(e)}

    async def process_single_file(
        self, 
        file_path: Union[str, Path],
        output_dir: Optional[str] = None
    ) -> ToolResponse:
        """Process a single file
        
        Args:
            file_path: Path to file to process
            output_dir: Optional override for output directory
            
        Returns:
            Tuple[bool, Dict]: (success, result_data)
            - success: Whether processing completed successfully
            - result_data: Dictionary containing processed article data or error
        """
        try:
            if output_dir:
                self.config.output_dir = output_dir
                
            file_path = Path(file_path)
            result = await self.process_file(file_path)
            
            if result:
                return True, {
                    "file": str(file_path),
                    "title": result.title,
                    "tags": result.tags,
                    "output_file": str(Path(self.config.output_dir) / f"{result.title}.md")
                }
            else:
                return False, {"error": f"Failed to process {file_path}"}

        except Exception as e:
            logger.error(f"Error processing file {file_path}: {str(e)}")
            return False, {"error": str(e)}

    async def process_file_with_semaphore(self, file_path: Path) -> Optional[KnowledgeArticle]:
        """Process a single file with concurrency control"""
        async with self.semaphore:
            return await self.process_file(file_path)

    async def process_file(self, file_path: Path) -> Optional[KnowledgeArticle]:
        """Main processing pipeline for a single file"""
        try:
            # Extract content from file
            content = await self._read_file(file_path)
            
            # Analyze and enhance content
            analysis = await self.content_analyzer.analyze_content(content, str(file_path))
            
            # Find related content
            related_articles = await self.kb_manager.find_related_articles(content)
            
            # Flatten hierarchical tags into a list
            all_tags = []
            if isinstance(analysis.get('tags'), dict):
                for rarity, tags in analysis['tags'].items():
                    all_tags.extend(tags)
            else:
                all_tags = analysis.get('tags', [])
            
            # Create knowledge article
            article = KnowledgeArticle(
                title=analysis["title"],
                content=analysis["formatted_content"],
                tags=all_tags,
                references=related_articles,
                source_file=str(file_path),
                metadata={"tags": analysis.get("tags", {})}  # Store original hierarchical tags in metadata
            )
            
            # Store in knowledge base
            await self.kb_manager.add_article(article)
            
            # Save as markdown file
            file_name = analysis.get("file_name", article.title.lower().replace(' ', '-'))
            output_path = Path(self.config.output_dir) / f"{file_name}.md"
            await self._save_markdown(article, output_path)
            
            return article
            
        except Exception as e:
            logger.error(f"Error processing {file_path}: {str(e)}")
            return None

    async def _read_file(self, file_path: Path) -> str:
        """Route file to appropriate reader based on extension"""
        suffix = file_path.suffix.lower()
        if suffix == '.pdf':
            return await self.document_reader.read_pdf(file_path)
        elif suffix in ['.md', '.txt']:
            return await self.document_reader.read_markdown(file_path)
        elif suffix == '.docx':
            return await self.document_reader.read_docx(file_path)
        elif suffix in ['.csv', '.xlsx']:
            return await self.document_reader.read_spreadsheet(file_path)
        elif suffix == '.pptx':
            return await self.document_reader.read_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {suffix}")

    async def _save_markdown(self, article: KnowledgeArticle, output_path: Path):
        """Generate and save markdown file with consistent formatting"""
        # Extract main heading from content, prioritizing first h1
        content_lines = article.content.split('\n')
        
        # Find all h1 headings
        h1_headings = [line.strip('# ') for line in content_lines 
                      if line.strip().startswith('# ') and not line.strip().startswith('## ')]
        
        # Use first h1 as main title, convert others to h2
        main_title = h1_headings[0] if h1_headings else article.title.replace('#', '').strip()
        content = article.content
        
        # Convert subsequent h1s to h2s and ensure proper spacing
        if len(h1_headings) > 1:
            for heading in h1_headings[1:]:
                content = content.replace(f'# {heading}', f'## {heading}')
        
        # Ensure proper spacing around headings and lists
        content = content.replace('\n#', '\n\n#')  # Add space before headings
        content = content.replace('\n-', '\n\n-')  # Add space before list items
        content = content.replace('\n\n\n', '\n\n')  # Remove extra blank lines
        
        # Format tags with color coding by rarity
        def format_tag_section(tags: List[str], color: str, rarity: str) -> str:
            if not tags:
                return ""
            formatted_tags = [
                (tag if tag.startswith('#') else f'#{tag}').lower().replace(' ', '-')
                for tag in tags
            ]
            # Map rarity levels to colored emojis
            emoji_map = {
                'white': '⚪',   # Common
                'green': '🟢',   # Fine
                'blue': '🔵',    # Rare
                'purple': '🟣',  # Epic
                'gold': '🟡'     # Legendary
            }
            emoji = emoji_map.get(color, '⚪')
            return f"{emoji} {' '.join(formatted_tags)}\n"  # Removed the "# {rarity}" label
        
        # Combine all tag sections
        tag_sections = []
        if article.metadata.get('tags'):
            tags = article.metadata['tags']
            tag_sections.extend([
                format_tag_section(tags.get('common', []), 'white', 'Common'),
                format_tag_section(tags.get('fine', []), 'green', 'Fine'),
                format_tag_section(tags.get('rare', []), 'blue', 'Rare'),
                format_tag_section(tags.get('epic', []), 'purple', 'Epic'),
                format_tag_section(tags.get('legendary', []), 'gold', 'Legendary')
            ])
        
        formatted_tags = ''.join(tag_sections) or ' '.join(
            f"⚪ {(tag if tag.startswith('#') else f'#{tag}').lower().replace(' ', '-')}"
            for tag in article.tags
        )
        
        # Generate markdown with consistent template and proper spacing
        template = f"""Title: 📎 {main_title}
Created: ⏰ {article.created_date.strftime('%Y-%m-%d')} {article.created_date.strftime('%H:%m')}
Tags: 🏷️
{formatted_tags}

---

## 📝 Knowledge Article

{content.strip()}

---

## 🔗 References

### 🌐 Connected Notes

{chr(10).join([f'- [[{ref}]]' for ref in article.references]) if article.references else '- No connected notes yet'}

### 📚 Sources

- Original file: {article.source_file}

---"""
        
        # Save with UTF-8 encoding
        output_path.write_text(template, encoding='utf-8')

# Convenience functions for direct usage
async def process_files(
    input_dir: str = "RAG_init",
    output_dir: str = "RAG_refined"
) -> ToolResponse:
    """Process all files in a directory (convenience function)"""
    factory = await MarkdownFactory.create(input_dir, output_dir)
    return await factory.process_directory()

async def process_file(
    file_path: Union[str, Path],
    output_dir: str = "RAG_refined"
) -> ToolResponse:
    """Process a single file (convenience function)"""
    factory = await MarkdownFactory.create(output_dir=output_dir)
    return await factory.process_single_file(file_path)

def main():
    """CLI entry point"""
    parser = argparse.ArgumentParser(
        description="Process documents into markdown knowledge articles"
    )
    parser.add_argument(
        "--input-dir",
        default="RAG_init",
        help="Directory containing source files"
    )
    parser.add_argument(
        "--output-dir",
        default="RAG_refined",
        help="Directory for processed markdown files"
    )
    args = parser.parse_args()

    # Run async process
    success, result = asyncio.run(process_files(args.input_dir, args.output_dir))
    
    if success:
        print("\n✅ Processing complete!")
        print(f"Processed: {result['stats']['successful']} files")
        print(f"Failed: {result['stats']['failed']} files")
        if result['failed_files']:
            print("\nFailed files:")
            for f in result['failed_files']:
                print(f"- {f}")
    else:
        print("\n❌ Processing failed!")
        print(f"Error: {result.get('error', 'Unknown error')}")

if __name__ == "__main__":
    main()
